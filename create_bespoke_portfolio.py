import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Palette: Biru Putih Minimalis Modern Profesional
    BG_WHITE = RGBColor(255, 255, 255)
    BG_LIGHT_SLATE = RGBColor(248, 250, 252) # #F8FAFC
    BG_CARD_WHITE = RGBColor(255, 255, 255)
    
    NAVY_PRIMARY = RGBColor(15, 36, 55)      # #0F2437 - Deep Executive Navy
    NAVY_SECONDARY = RGBColor(30, 58, 138)   # #1E3A8A - Rich Navy
    BLUE_ACCENT = RGBColor(37, 99, 235)      # #2563EB - Electric Cobalt
    CYAN_ACCENT = RGBColor(14, 165, 233)     # #0EA5E9 - Sky Blue
    
    TEXT_MAIN = RGBColor(15, 23, 42)         # #0F172A - Deep Charcoal
    TEXT_MUTED = RGBColor(71, 85, 105)       # #475569 - Slate
    TEXT_LIGHT = RGBColor(148, 163, 184)     # #94A3B8 - Light Slate
    TEXT_WHITE = RGBColor(255, 255, 255)
    
    BORDER_SUBTLE = RGBColor(226, 232, 240)  # #E2E8F0 - Clean Card Border
    BORDER_ACCENT = RGBColor(191, 219, 254)  # #BFDBFE - Soft Blue Border
    PILL_BG_BLUE = RGBColor(239, 246, 255)   # #EFF6FF - Soft Blue Pill

    FONT_FAMILY = "Segoe UI"

    def set_background(slide, color=BG_LIGHT_SLATE):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_card(slide, left, top, width, height, fill_color=BG_CARD_WHITE, border_color=BORDER_SUBTLE, rounded=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        card = slide.shapes.add_shape(shape_type, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1)
        else:
            card.line.fill.background()
        return card

    def add_pill_tag(slide, left, top, width, height, text, bg_color=PILL_BG_BLUE, text_color=BLUE_ACCENT, font_size=10):
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        pill.fill.solid()
        pill.fill.fore_color.rgb = bg_color
        pill.line.color.rgb = BORDER_ACCENT
        pill.line.width = Pt(0.75)
        tf = pill.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_FAMILY
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = text_color
        return pill

    def add_header(slide, tag_text, title_text, subtitle_text=""):
        # Tag
        add_pill_tag(slide, Inches(0.8), Inches(0.5), Inches(3.6), Inches(0.32), tag_text, PILL_BG_BLUE, BLUE_ACCENT, 9)
        # Title
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.88), Inches(11.7), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = NAVY_PRIMARY

        # Subtitle
        if subtitle_text:
            tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.35))
            tf_sub = tb_sub.text_frame
            tf_sub.word_wrap = True
            tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_text
            p_sub.font.name = FONT_FAMILY
            p_sub.font.size = Pt(12)
            p_sub.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 1: COVER SLIDE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1, BG_LIGHT_SLATE)

    # Decorative top navy bar
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333333), Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = BLUE_ACCENT
    top_bar.line.fill.background()

    # Left Column: Hero Information Card
    left_hero = add_card(slide1, Inches(0.8), Inches(0.65), Inches(7.5), Inches(6.2), BG_CARD_WHITE, BORDER_SUBTLE)
    
    # Tag
    add_pill_tag(slide1, Inches(1.2), Inches(1.0), Inches(3.2), Inches(0.34), "PORTOFOLIO PROFESIONAL 2026", PILL_BG_BLUE, BLUE_ACCENT, 9)

    # Name
    name_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.45), Inches(6.8), Inches(0.85))
    tf = name_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    p.text = "Khairul Raihan Hidayat, S.Kom."
    p.font.name = FONT_FAMILY
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    # Track / Role
    role_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.35), Inches(6.8), Inches(0.45))
    tf = role_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    p.text = "Data Analyst & Agentic AI Specialist"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT

    # Academic Distinction Badge
    honor_pill = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.9), Inches(5.8), Inches(0.4))
    honor_pill.fill.solid()
    honor_pill.fill.fore_color.rgb = RGBColor(254, 243, 199) # Warm Gold/Amber tint
    honor_pill.line.color.rgb = RGBColor(245, 158, 11)
    honor_pill.line.width = Pt(1)
    tf = honor_pill.text_frame
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "⭐ Sarjana Sistem Informasi (Data Science) | IPK 3.88 / 4.00 (Magna Cum Laude)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(146, 64, 14)

    # Bio Summary
    bio_box = slide1.shapes.add_textbox(Inches(1.2), Inches(3.45), Inches(6.7), Inches(1.3))
    tf = bio_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    p.text = "Lulusan terbaik Universitas Budi Luhur yang memadukan keahlian analisis data kuantitatif, pemodelan Machine Learning (NLP/Klasifikasi), visualisasi data eksekutif (Tableau & SQL), serta rekayasa otomasi alur kerja modern berbasis Google Antigravity Agentic AI."
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11.5)
    p.font.color.rgb = TEXT_MUTED

    # Contact Cards Grid (2x2)
    contacts = [
        ("📧 Email", "khairulraihan1204@gmail.com"),
        ("📱 WhatsApp", "+62 822-1081-3652"),
        ("💼 LinkedIn", "linkedin.com/in/khairulraihan"),
        ("🌐 Web Portfolio", "khairulraihan.github.io")
    ]
    for idx, (label, val) in enumerate(contacts):
        c_left = Inches(1.2 + (idx % 2) * 3.4)
        c_top = Inches(4.9 + (idx // 2) * 0.8)
        card = add_card(slide1, c_left, c_top, Inches(3.25), Inches(0.7), BG_LIGHT_SLATE, BORDER_SUBTLE)
        tb = slide1.shapes.add_textbox(c_left + Inches(0.12), c_top + Inches(0.08), Inches(3.0), Inches(0.55))
        tf = tb.text_frame
        tf.margin_left = tf.margin_top = 0
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(9.5)
        p1.font.bold = True
        p1.font.color.rgb = BLUE_ACCENT
        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_MAIN

    # Right Column: Visual Portrait & Key Metrics Card
    right_card = add_card(slide1, Inches(8.6), Inches(0.65), Inches(3.9), Inches(6.2), NAVY_PRIMARY, border_color=None)
    
    # Add profile photo if exists
    profile_path = os.path.abspath("assets/images/profile.jpg")
    if os.path.exists(profile_path):
        slide1.shapes.add_picture(profile_path, Inches(9.3), Inches(1.1), Inches(2.5), Inches(2.5))
        # Photo Border Accent Frame
        frame = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.25), Inches(1.05), Inches(2.6), Inches(2.6))
        frame.fill.background()
        frame.line.color.rgb = CYAN_ACCENT
        frame.line.width = Pt(2)

    # Text below photo inside Navy Card
    stats_box = slide1.shapes.add_textbox(Inches(8.8), Inches(3.8), Inches(3.5), Inches(2.8))
    tf = stats_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    
    p = tf.paragraphs[0]
    p.text = "METRIK & PENCAPAIAN UTAMA"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    metrics_list = [
        ("3.88 / 4.00", "IPK Lulusan Terbaik (Magna Cum Laude)"),
        ("88.41%", "Akurasi Model Linear SVM pada Skripsi NLP"),
        ("8.950+", "Catatan Karyawan Dianalisis di Tableau BI"),
        ("~70%", "Efisiensi Otomasi Workflow Agentic AI")
    ]
    for val, desc in metrics_list:
        p_val = tf.add_paragraph()
        p_val.text = f"• {val} — {desc}"
        p_val.font.name = FONT_FAMILY
        p_val.font.size = Pt(10)
        p_val.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 2: PROFIL PROFESIONAL & 3 PILAR STRATEGIS
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2, BG_LIGHT_SLATE)
    add_header(slide2, "EXECUTIVE PROFILE", "Profil Profesional & 3 Pilar Kompetensi Utama", "Mengintegrasikan ketajaman analitis, kecakapan business intelligence, dan orkestrasi Agentic AI.")

    # Executive Summary Banner
    summary_card = add_card(slide2, Inches(0.8), Inches(1.95), Inches(11.7), Inches(1.1), BG_CARD_WHITE, BORDER_SUBTLE)
    tb = slide2.shapes.add_textbox(Inches(1.0), Inches(2.05), Inches(11.3), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    p.text = "Khairul Raihan Hidayat adalah Sarjana Sistem Informasi (Data Science) dari Universitas Budi Luhur dengan predikat kelulusan Magna Cum Laude (IPK 3.88). Memiliki fondasi analitika data yang kuat, spesialisasi dalam pemrosesan bahasa alami (NLP), pembuatan interactive executive dashboards menggunakan Tableau & SQL, serta pelopor implementasi otomasi alur kerja digital berbasis sistem multi-agen Google Antigravity."
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11.5)
    p.font.color.rgb = TEXT_MAIN

    # 3 Strategic Pillar Cards
    pillars = [
        {
            "tag": "PILAR 01",
            "title": "Machine Learning & NLP Terapan",
            "subtitle": "Predictive Modeling & Text Analytics",
            "points": [
                "Pemodelan klasifikasi terapan: Support Vector Machine (SVM), Naive Bayes, K-Nearest Neighbors (KNN).",
                "Pipeline text preprocessing komprehensif: Case folding, normalisasi slang/kata gaul, stopword removal Sastrawi.",
                "Ekstraksi fitur Sparse Matrix berdimensi tinggi menggunakan CountVectorizer & TF-IDF (11.790 vocab unik).",
                "Evaluasi ketat via Stratified Split, Confusion Matrix, Precision, Recall, dan F1-Score."
            ],
            "accent": BLUE_ACCENT
        },
        {
            "tag": "PILAR 02",
            "title": "Business Intelligence & EDA",
            "subtitle": "Interactive Dashboards & Decision Support",
            "points": [
                "Perancangan visualisasi data interaktif dan storytelling eksekutif menggunakan Tableau Desktop.",
                "Eksplorasi dataset korporat skala ribuan baris (HR workforce 8.950 records, e-commerce retail).",
                "Query data relasional kompleks dengan SQL: Aggregations, CTEs, Joins, Window Functions.",
                "Analisis metrik bisnis krusial: Retention Rate, Attrition Risk, Regional Profit Margin, dan Pola Pembelian."
            ],
            "accent": NAVY_SECONDARY
        },
        {
            "tag": "PILAR 03",
            "title": "Agentic AI & Workflow Automation",
            "subtitle": "Multi-Agent Systems & Process Engineering",
            "points": [
                "Perancangan arsitektur multi-agen cerdas memanfaatkan ekosistem Google Antigravity Framework.",
                "Otomasi rantai produksi konten dari kurasi transkrip podcast hingga short-form video viral 9:16.",
                "Integrasi AI vision, programmatic scripting (CapCut/FFmpeg), dan prompt engineering presisi.",
                "Mencapai ~70% efisiensi waktu turnaround produksi dan eliminasi bottleneck manual."
            ],
            "accent": CYAN_ACCENT
        }
    ]

    for idx, col in enumerate(pillars):
        c_left = Inches(0.8 + idx * 4.0)
        c_top = Inches(3.25)
        c_w = Inches(3.7)
        c_h = Inches(3.7)

        # Card shape
        card = add_card(slide2, c_left, c_top, c_w, c_h, BG_CARD_WHITE, BORDER_SUBTLE)

        # Top Accent Strip
        strip = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, c_top, c_w, Inches(0.12))
        strip.fill.solid()
        strip.fill.fore_color.rgb = col["accent"]
        strip.line.fill.background()

        # Content inside card
        tb = slide2.shapes.add_textbox(c_left + Inches(0.25), c_top + Inches(0.25), c_w - Inches(0.5), c_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = 0

        p0 = tf.paragraphs[0]
        p0.text = col["tag"]
        p0.font.name = FONT_FAMILY
        p0.font.size = Pt(9.5)
        p0.font.bold = True
        p0.font.color.rgb = col["accent"]

        p1 = tf.add_paragraph()
        p1.text = col["title"]
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_PRIMARY

        p2 = tf.add_paragraph()
        p2.text = col["subtitle"]
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(10)
        p2.font.italic = True
        p2.font.color.rgb = TEXT_MUTED

        for pt in col["points"]:
            p_pt = tf.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.name = FONT_FAMILY
            p_pt.font.size = Pt(9.5)
            p_pt.font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 3: KEAHLIAN TEKNIS & EKOSISTEM TOOLING
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3, BG_LIGHT_SLATE)
    add_header(slide3, "TECHNICAL TAXONOMY", "Keahlian Teknis & Ekosistem Tooling", "Kapabilitas teknis end-to-end dari data ingestion, preprocessing, pemodelan, hingga deployment.")

    skill_groups = [
        {
            "category": "Bahasa Pemrograman & Querying",
            "skills": [
                ("Python", "Pandas, NumPy, Scikit-Learn, SciPy, Matplotlib, Seaborn"),
                ("SQL (Structured Query)", "PostgreSQL, MySQL, Window Functions, Complex Joins, CTEs"),
                ("R Language", "Statistical computing & exploratory baseline modeling"),
                ("Data Manipulation", "Regular Expressions (Regex), JSON, CSV, API Ingestion")
            ]
        },
        {
            "category": "Data Science & Machine Learning",
            "skills": [
                ("Algoritma Klasifikasi", "Support Vector Machine (Linear SVM), Naive Bayes, KNN"),
                ("Natural Language Processing", "CountVectorizer, TF-IDF, Sastrawi Stemmer, Lexicon-based"),
                ("Evaluasi Model", "Stratified K-Fold, Confusion Matrix, F1-Score, ROC-AUC"),
                ("Hyperparameter Tuning", "GridSearchCV, RandomSearch, Pipeline Automation")
            ]
        },
        {
            "category": "Business Intelligence & Visualization",
            "skills": [
                ("Tableau Desktop & Public", "Calculated Fields, LOD Expressions, Parameter Actions, Sets"),
                ("Web Application", "Streamlit (Live interactive Machine Learning demo web apps)"),
                ("Executive Storytelling", "KPI Dashboarding, Geospatial Mapping, Scatter & Funnel Plots"),
                ("Exploratory Data Analysis", "Trend detection, outlier identification, cohort retention")
            ]
        },
        {
            "category": "Agentic AI, DevOps & Tooling",
            "skills": [
                ("Agentic AI Frameworks", "Google Antigravity, Multi-agent autonomous workflow design"),
                ("Version Control", "Git, GitHub, Branching Strategy, CI/CD automation"),
                ("Workflow Automation", "Python Automation Scripts, CapCut API batching, FFmpeg"),
                ("Development Environments", "VS Code, Jupyter Notebook, Google Colab, Windows PowerShell")
            ]
        }
    ]

    for idx, grp in enumerate(skill_groups):
        r = idx // 2
        c = idx % 2
        c_left = Inches(0.8 + c * 6.0)
        c_top = Inches(2.0 + r * 2.55)
        c_w = Inches(5.7)
        c_h = Inches(2.35)

        card = add_card(slide3, c_left, c_top, c_w, c_h, BG_CARD_WHITE, BORDER_SUBTLE)
        
        # Header strip inside card
        header_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, c_top, c_w, Inches(0.45))
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = PILL_BG_BLUE
        header_bg.line.color.rgb = BORDER_ACCENT
        tf_h = header_bg.text_frame
        tf_h.margin_left = Inches(0.15)
        p_h = tf_h.paragraphs[0]
        p_h.text = f"❖  {grp['category']}"
        p_h.font.name = FONT_FAMILY
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_SECONDARY

        # Content items
        tb = slide3.shapes.add_textbox(c_left + Inches(0.2), c_top + Inches(0.55), c_w - Inches(0.4), c_h - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = 0

        for i, (title, details) in enumerate(grp["skills"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {title}: "
            p.font.name = FONT_FAMILY
            p.font.size = Pt(9.5)
            p.font.bold = True
            p.font.color.rgb = NAVY_PRIMARY

            # add detail run
            r_det = p.add_run()
            r_det.text = details
            r_det.font.bold = False
            r_det.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 4: PROJECT 1 - SKRIPSI NLP (OVERVIEW & MODEL COMPARISON)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4, BG_LIGHT_SLATE)
    add_header(slide4, "PROJECT 01: MACHINE LEARNING SKRIPSI", "Komparasi SVM, Naive Bayes & KNN: Analisis Sentimen YouTube", "Studi Kasus Respons Netizen terhadap Pengungkapan Kasus Korupsi Videografer di Sumatera Utara (Kanal Ferry Irwandi).")

    # Left Column: Project Overview & Scientific Findings
    left_card = add_card(slide4, Inches(0.8), Inches(1.95), Inches(5.9), Inches(5.0), BG_CARD_WHITE, BORDER_SUBTLE)
    
    tb_l = slide4.shapes.add_textbox(Inches(1.05), Inches(2.1), Inches(5.4), Inches(4.7))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_top = 0

    p = tf_l.paragraphs[0]
    p.text = "RINGKASAN PENELITIAN & METODOLOGI"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT

    p_body = tf_l.add_paragraph()
    p_body.text = "Penelitian skripsi ini membandingkan kinerja tiga algoritma Machine Learning terpopuler dalam mengklasifikasikan respons sentimen publik pada 9.269 komentar YouTube kasus korupsi. Dataset dikonversi menjadi representasi numerik menggunakan CountVectorizer menghasilkan ruang sparse matrix berdimensi 11.790 fitur unik."
    p_body.font.name = FONT_FAMILY
    p_body.font.size = Pt(10)
    p_body.font.color.rgb = TEXT_MAIN

    # Stats Row
    p_stats_head = tf_l.add_paragraph()
    p_stats_head.text = "TEMUAN UTAMA KOMPARASI ALGORITMA:"
    p_stats_head.font.name = FONT_FAMILY
    p_stats_head.font.size = Pt(10.5)
    p_stats_head.font.bold = True
    p_stats_head.font.color.rgb = NAVY_PRIMARY

    findings = [
        ("Linear SVM (88.41% Akurasi)", "Mendominasi secara mutlak di seluruh skenario split (60:40, 70:30, 80:20). Sangat tangguh mengatasi fenomena high-dimensional sparsity."),
        ("Multinomial Naive Bayes (81.03%)", "Performa klasifikasi moderat. Cepat dalam inferensi, namun terkendala asumsi independensi bersyarat pada konteks kalimat sarkasme/kritik publik."),
        ("K-Nearest Neighbors (68.22%)", "Performa terendah akibat kutukan dimensi (curse of dimensionality), di mana perhitungan jarak Euclidean kehilangan diskriminasi pada 11.790 fitur.")
    ]
    for m_title, m_desc in findings:
        p_m = tf_l.add_paragraph()
        p_m.text = f"✓ {m_title}: "
        p_m.font.name = FONT_FAMILY
        p_m.font.size = Pt(9.5)
        p_m.font.bold = True
        p_m.font.color.rgb = BLUE_ACCENT
        run = p_m.add_run()
        run.text = m_desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    # Right Column: Embed Thesis Accuracy Comparison Slide (Slide 11)
    right_card = add_card(slide4, Inches(7.0), Inches(1.95), Inches(5.5), Inches(5.0), BG_CARD_WHITE, BORDER_SUBTLE)
    
    thesis_img_11 = os.path.abspath("skripsi_slide_preview/slide_11.png")
    if os.path.exists(thesis_img_11):
        slide4.shapes.add_picture(thesis_img_11, Inches(7.15), Inches(2.1), Inches(5.2), Inches(3.8))
        
        # Caption below picture
        tb_cap = slide4.shapes.add_textbox(Inches(7.15), Inches(6.0), Inches(5.2), Inches(0.8))
        tf_cap = tb_cap.text_frame
        tf_cap.word_wrap = True
        tf_cap.margin_left = tf_cap.margin_top = 0
        p_cap = tf_cap.paragraphs[0]
        p_cap.text = "Grafik Komparasi Resmi Sidang Skripsi (Universitas Budi Luhur, 2026): Linear SVM mencapai puncak akurasi 88.41% pada rasio split 80:20."
        p_cap.font.name = FONT_FAMILY
        p_cap.font.size = Pt(9.5)
        p_cap.font.italic = True
        p_cap.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 5: PROJECT 1 - PIPELINE NLP, CONFUSION MATRIX & STREAMLIT APP
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5, BG_LIGHT_SLATE)
    add_header(slide5, "PROJECT 01: ARCHITECTURE & EVALUATION", "Arsitektur Pipeline NLP & Analisis Confusion Matrix", "Dari ekstraksi YouTube Data API v3, kamus slang normalisasi, hingga deploy Streamlit Web App.")

    # Top Pipeline Cards (5 stages)
    pipeline_steps = [
        ("Tahap 1: Ingestion", "9.269 Komentar Mentah diekstrak via YouTube Data API v3 Google Cloud."),
        ("Tahap 2: Cleaning", "Pembersihan URL, hashtag, angka, tanda baca & emoji via Regex."),
        ("Tahap 3: Normalisasi", "Mapping kata gaul/slang bahasa Indonesia ke bentuk baku terstruktur."),
        ("Tahap 4: Stopwords", "Penghapusan noise kata sambung korpus Sastrawi + kamus buatan."),
        ("Tahap 5: Feature Vector", "CountVectorizer sparse matrix menghasilkan 11.790 vocab unik.")
    ]

    for idx, (p_title, p_desc) in enumerate(pipeline_steps):
        p_left = Inches(0.8 + idx * 2.4)
        p_top = Inches(1.95)
        p_w = Inches(2.25)
        p_h = Inches(1.25)

        card = add_card(slide5, p_left, p_top, p_w, p_h, BG_CARD_WHITE, BORDER_SUBTLE)
        tb = slide5.shapes.add_textbox(p_left + Inches(0.12), p_top + Inches(0.1), p_w - Inches(0.24), p_h - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = 0
        p0 = tf.paragraphs[0]
        p0.text = p_title
        p0.font.name = FONT_FAMILY
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = BLUE_ACCENT
        p1 = tf.add_paragraph()
        p1.text = p_desc
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(8.5)
        p1.font.color.rgb = TEXT_MUTED

    # Bottom Left: Embed Confusion Matrix Image (Slide 13)
    bot_left_card = add_card(slide5, Inches(0.8), Inches(3.35), Inches(5.9), Inches(3.6), BG_CARD_WHITE, BORDER_SUBTLE)
    thesis_img_13 = os.path.abspath("skripsi_slide_preview/slide_13.png")
    if os.path.exists(thesis_img_13):
        slide5.shapes.add_picture(thesis_img_13, Inches(0.95), Inches(3.45), Inches(5.6), Inches(3.35))

    # Bottom Right: Rigorous Statistical Breakdown & Streamlit Web App
    bot_right_card = add_card(slide5, Inches(7.0), Inches(3.35), Inches(5.5), Inches(3.6), BG_CARD_WHITE, BORDER_SUBTLE)
    tb_r = slide5.shapes.add_textbox(Inches(7.25), Inches(3.5), Inches(5.0), Inches(3.3))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_top = 0

    p_rt = tf_r.paragraphs[0]
    p_rt.text = "ANALISIS VALIDASI MATRIKS & DEPLOYMENT"
    p_rt.font.name = FONT_FAMILY
    p_rt.font.size = Pt(11)
    p_rt.font.bold = True
    p_rt.font.color.rgb = NAVY_PRIMARY

    breakdown_points = [
        ("Total Data Uji Evaluasi", "1.476 Baris komentar (rasio split 80:20)"),
        ("True Positive & True Negative", "808 Sentimen Positif & 497 Sentimen Negatif terklasifikasi akurat."),
        ("Presisi Kelas Negatif (88.6%)", "Model sangat selektif dan tidak gegabah melabeli kritik masyarakat sebagai ujaran negatif tanpa konteks."),
        ("Recall Kelas Positif (92.7%)", "Model sangat sensitif dan akurat dalam menangkap narasi dukungan terhadap penegakan hukum korupsi."),
        ("Streamlit Interactive Deployment", "Model diekspor via Pickle dan diintegrasikan ke web app Streamlit untuk inferensi sentimen real-time.")
    ]

    for b_title, b_desc in breakdown_points:
        p_b = tf_r.add_paragraph()
        p_b.text = f"• {b_title}: "
        p_b.font.name = FONT_FAMILY
        p_b.font.size = Pt(9.5)
        p_b.font.bold = True
        p_b.font.color.rgb = BLUE_ACCENT
        run = p_b.add_run()
        run.text = b_desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 6: PROJECT 2 - TABLEAU HR ANALYTICS DASHBOARD
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6, BG_LIGHT_SLATE)
    add_header(slide6, "PROJECT 02: BUSINESS INTELLIGENCE", "Tableau HR Analytics: Workforce & Retention Dashboard", "Analisis komprehensif 8.950 catatan karyawan untuk mendeteksi tren retensi, kompensasi, dan performa.")

    # Left Column: Embed Real Tableau HR Dashboard Screenshot
    left_hr_card = add_card(slide6, Inches(0.8), Inches(1.95), Inches(6.8), Inches(5.0), BG_CARD_WHITE, BORDER_SUBTLE)
    hr_img = os.path.abspath("assets/images/project-hr-analytics.png")
    if os.path.exists(hr_img):
        slide6.shapes.add_picture(hr_img, Inches(0.95), Inches(2.1), Inches(6.5), Inches(4.3))
        tb_hr_cap = slide6.shapes.add_textbox(Inches(0.95), Inches(6.45), Inches(6.5), Inches(0.4))
        tf_cap = tb_hr_cap.text_frame
        tf_cap.margin_left = tf_cap.margin_top = 0
        p_c = tf_cap.paragraphs[0]
        p_c.text = "Interactive Tableau Desktop Dashboard: Employee Retention, State Map & Performance Scatter Plot."
        p_c.font.name = FONT_FAMILY
        p_c.font.size = Pt(9)
        p_c.font.italic = True
        p_c.font.color.rgb = TEXT_MUTED

    # Right Column: Metrics & Business Storytelling
    right_hr_card = add_card(slide6, Inches(7.8), Inches(1.95), Inches(4.7), Inches(5.0), BG_CARD_WHITE, BORDER_SUBTLE)
    
    # 3 Metric KPI Cards inside right card
    kpis = [
        ("8.950", "TOTAL EMPLOYEES", "Basis data tenaga kerja dianalisis"),
        ("89.2%", "RETENTION RATE", "Tingkat loyalitas karyawan aktif"),
        ("10.8%", "TURNOVER RATE", "Tingkat atrisi teridentifikasi")
    ]
    for idx, (val, title, sub) in enumerate(kpis):
        k_top = Inches(2.1 + idx * 0.95)
        k_card = add_card(slide6, Inches(8.0), k_top, Inches(4.3), Inches(0.85), PILL_BG_BLUE, BORDER_ACCENT)
        tb_k = slide6.shapes.add_textbox(Inches(8.15), k_top + Inches(0.08), Inches(4.0), Inches(0.7))
        tf_k = tb_k.text_frame
        tf_k.margin_left = tf_k.margin_top = 0
        p_val = tf_k.paragraphs[0]
        p_val.text = f"{val}  "
        p_val.font.name = FONT_FAMILY
        p_val.font.size = Pt(16)
        p_val.font.bold = True
        p_val.font.color.rgb = NAVY_PRIMARY
        
        r_t = p_val.add_run()
        r_t.text = f"— {title}"
        r_t.font.size = Pt(10)
        r_t.font.bold = True
        r_t.font.color.rgb = BLUE_ACCENT
        
        p_s = tf_k.add_paragraph()
        p_s.text = sub
        p_s.font.name = FONT_FAMILY
        p_s.font.size = Pt(8.5)
        p_s.font.color.rgb = TEXT_MUTED

    # Insights Text Box below KPIs
    tb_hr_ins = slide6.shapes.add_textbox(Inches(8.0), Inches(5.05), Inches(4.3), Inches(1.8))
    tf_ins = tb_hr_ins.text_frame
    tf_ins.word_wrap = True
    tf_ins.margin_left = tf_ins.margin_top = 0

    p_h = tf_ins.paragraphs[0]
    p_h.text = "STRATEGIC HR INSIGHTS:"
    p_h.font.name = FONT_FAMILY
    p_h.font.size = Pt(10)
    p_h.font.bold = True
    p_h.font.color.rgb = NAVY_PRIMARY

    hr_points = [
        "Kompensasi vs Performa: Analisis scatter plot mengidentifikasi retensi talenta berkinerja tinggi rentan terhadap stagnasi gaji.",
        "Geographic Mapping: Visualisasi sebaran kantor cabang membantu pemerataan beban kerja dan perencanaan tunjangan regional.",
        "Rekomendasi Kebijakan: Penyesuaian skema review performa semesteran untuk memitigasi turnover pada divisi berisiko tinggi."
    ]
    for pt in hr_points:
        p_pt = tf_ins.add_paragraph()
        p_pt.text = f"• {pt}"
        p_pt.font.name = FONT_FAMILY
        p_pt.font.size = Pt(8.5)
        p_pt.font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 7: PROJECT 3 - AGENTIC AI VIDEO & AUTOMATION PIPELINE
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7, BG_LIGHT_SLATE)
    add_header(slide7, "PROJECT 03: AGENTIC AI AUTOMATION", "Autonomous Multi-Agent AI Video Pipeline", "Orkestrasi agen otonom dengan Google Antigravity untuk kurasi & rendering video pendek berkecepatan tinggi.")

    # Top Metrics Banner (3 Cards)
    metric_cols = [
        ("~70%", "REDUKSI TURNAROUND", "Penyusutan waktu dari transkrip mentah ke video siap tayang."),
        ("3x", "VOLUME PRODUKSI", "Peningkatan output klip pendek harian multi-platform."),
        ("100%", "KONSISTENSI BRANDING", "Automated subtitling, aspect ratio 9:16, and b-roll placement.")
    ]
    for idx, (m_val, m_t, m_d) in enumerate(metric_cols):
        m_left = Inches(0.8 + idx * 4.0)
        m_card = add_card(slide7, m_left, Inches(1.95), Inches(3.7), Inches(1.0), BG_CARD_WHITE, BORDER_SUBTLE)
        tb_m = slide7.shapes.add_textbox(m_left + Inches(0.15), Inches(2.05), Inches(3.4), Inches(0.8))
        tf_m = tb_m.text_frame
        tf_m.word_wrap = True
        tf_m.margin_left = tf_m.margin_top = 0
        p_v = tf_m.paragraphs[0]
        p_v.text = f"{m_val} "
        p_v.font.name = FONT_FAMILY
        p_v.font.size = Pt(18)
        p_v.font.bold = True
        p_v.font.color.rgb = BLUE_ACCENT
        r_t = p_v.add_run()
        r_t.text = f" {m_t}"
        r_t.font.size = Pt(10)
        r_t.font.bold = True
        r_t.font.color.rgb = NAVY_PRIMARY
        p_d = tf_m.add_paragraph()
        p_d.text = m_d
        p_d.font.name = FONT_FAMILY
        p_d.font.size = Pt(8.5)
        p_d.font.color.rgb = TEXT_MUTED

    # 3 Multi-Agent Workflow Pillars
    workflow_agents = [
        {
            "tier": "AGEN 01: INGESTION & VIRALITY ANALYZER",
            "role": "Content Intelligence & Hook Detection",
            "desc": "Mengekstrak transkrip audio podcast berdurasi panjang dengan timestamp presisi. Menganalisis densitas emosi, keyword virality, dan mengidentifikasi potensi 'golden moments' 30-60 detik.",
            "tech": "Python, Speech-to-Text API, Semantic Scoring"
        },
        {
            "tier": "AGEN 02: ADAPTATION & SCRIPT SYNTHESIS",
            "role": "Narrative Structuring & Pacing",
            "desc": "Mengolah transkrip mentah menjadi format vertical short-form (9:16). Menambahkan formula hook 3 detik pembuka, teks ringkas berdaya pikat tinggi, serta instruksi transisi visual.",
            "tech": "Google Antigravity Multi-agent, Structured Prompting"
        },
        {
            "tier": "AGEN 03: PROGRAMMATIC VIDEO ASSEMBLER",
            "role": "Automated Editing & Rendering Execution",
            "desc": "Mengeksekusi otomasi editing melalui script CapCut & FFmpeg. Menghasilkan automated dynamic subtitles, visual sound effects, auto-framing pembicara, dan ekspor batch resolusi tinggi.",
            "tech": "FFmpeg Engine, CapCut API / Script Automation"
        }
    ]

    for idx, ag in enumerate(workflow_agents):
        a_left = Inches(0.8 + idx * 4.0)
        a_top = Inches(3.15)
        a_w = Inches(3.7)
        a_h = Inches(3.8)

        card = add_card(slide7, a_left, a_top, a_w, a_h, BG_CARD_WHITE, BORDER_SUBTLE)
        
        # Top strip
        strip = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, a_left, a_top, a_w, Inches(0.08))
        strip.fill.solid()
        strip.fill.fore_color.rgb = BLUE_ACCENT if idx == 1 else (CYAN_ACCENT if idx == 0 else NAVY_PRIMARY)
        strip.line.fill.background()

        tb = slide7.shapes.add_textbox(a_left + Inches(0.2), a_top + Inches(0.2), a_w - Inches(0.4), a_h - Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = 0

        p0 = tf.paragraphs[0]
        p0.text = ag["tier"]
        p0.font.name = FONT_FAMILY
        p0.font.size = Pt(9.5)
        p0.font.bold = True
        p0.font.color.rgb = BLUE_ACCENT

        p1 = tf.add_paragraph()
        p1.text = ag["role"]
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_PRIMARY

        p2 = tf.add_paragraph()
        p2.text = ag["desc"]
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_MAIN

        p3 = tf.add_paragraph()
        p3.text = f"Stack: {ag['tech']}"
        p3.font.name = FONT_FAMILY
        p3.font.size = Pt(9)
        p3.font.italic = True
        p3.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 8: PROJECT 4 - AMAZON SALES & PROFITABILITY DASHBOARD
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8, BG_LIGHT_SLATE)
    add_header(slide8, "PROJECT 04: COMMERCIAL ANALYTICS", "Amazon Commercial Sales & Regional Profitability", "Eksplorasi data transaksi e-commerce multi-regional untuk optimalisasi marjin dan logistik.")

    # 3 Analytical Focus Cards
    analysis_cards = [
        {
            "tag": "ANALISIS 01",
            "title": "Profitabilitas Regional & Biaya Logistik",
            "points": [
                "Pemetaan pendapatan kotor vs net margin di seluruh regional pasar komersial.",
                "Mengidentifikasi anomali di mana volume transaksi tinggi pada state tertentu mengalami erosi profit akibat tingginya biaya ongkos kirim (shipping freight).",
                "Rekomendasi: Penyesuaian threshold minimum basket size untuk gratis ongkir di zona logistik jauh."
            ]
        },
        {
            "tag": "ANALISIS 02",
            "title": "Kategori Produk & Kontribusi Marjin",
            "points": [
                "Penerapan Analisis Pareto (80/20) pada kategori produk (Technology, Office Supplies, Furniture).",
                "Memisahkan produk kategori 'High-Volume High-Margin' dari produk 'Cash Trap' yang over-diskon namun menghasilkan profit negatif.",
                "Rekomendasi: Re-alokasi anggaran promosi ke sub-kategori bernilai tambah tinggi."
            ]
        },
        {
            "tag": "ANALISIS 03",
            "title": "Segmentasi Pelanggan (B2B vs B2C)",
            "points": [
                "Eksplorasi perilaku pembelian antara segmen Korporasi, Home Office, dan Konsumen Langsung.",
                "Menemukan segmen Korporat memiliki repeat order rate 38% lebih konsisten dengan margin stabil.",
                "Rekomendasi: Perancangan skema dedicated account support dan tiering harga B2B berbasis volume."
            ]
        }
    ]

    for idx, card_data in enumerate(analysis_cards):
        c_left = Inches(0.8 + idx * 4.0)
        c_top = Inches(1.95)
        c_w = Inches(3.7)
        c_h = Inches(3.5)

        card = add_card(slide8, c_left, c_top, c_w, c_h, BG_CARD_WHITE, BORDER_SUBTLE)
        tb = slide8.shapes.add_textbox(c_left + Inches(0.2), c_top + Inches(0.2), c_w - Inches(0.4), c_h - Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = 0

        p0 = tf.paragraphs[0]
        p0.text = card_data["tag"]
        p0.font.name = FONT_FAMILY
        p0.font.size = Pt(9.5)
        p0.font.bold = True
        p0.font.color.rgb = BLUE_ACCENT

        p1 = tf.add_paragraph()
        p1.text = card_data["title"]
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_PRIMARY

        for pt in card_data["points"]:
            p_pt = tf.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.name = FONT_FAMILY
            p_pt.font.size = Pt(9)
            p_pt.font.color.rgb = TEXT_MAIN

    # Bottom Banner: Actionable Business Value
    val_card = add_card(slide8, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.3), PILL_BG_BLUE, BORDER_ACCENT)
    tb_val = slide8.shapes.add_textbox(Inches(1.0), Inches(5.75), Inches(11.3), Inches(1.1))
    tf_val = tb_val.text_frame
    tf_val.word_wrap = True
    tf_val.margin_left = tf_val.margin_top = 0

    p_vh = tf_val.paragraphs[0]
    p_vh.text = "DAMPAK STRATEGIS BISNIS YANG DIHASILKAN:"
    p_vh.font.name = FONT_FAMILY
    p_vh.font.size = Pt(10.5)
    p_vh.font.bold = True
    p_vh.font.color.rgb = NAVY_PRIMARY

    p_vb = tf_val.add_paragraph()
    p_vb.text = "Dashboard Tableau ini mentransformasi ratusan ribu data transaksi mentah menjadi insight aksi komersial nyata. Manajemen dapat langsung mendeteksi kebocoran margin, memangkas diskon kontraproduktif, serta mengoptimalkan jalur pengiriman antar-negara bagian untuk memulihkan profitabilitas e-commerce secara terukur."
    p_vb.font.name = FONT_FAMILY
    p_vb.font.size = Pt(10)
    p_vb.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 9: PENDIDIKAN, PRESTASI & SERTIFIKASI PROFESIONAL
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9, BG_LIGHT_SLATE)
    add_header(slide9, "CREDENTIALS & HONORS", "Pendidikan, Prestasi Akademik & Sertifikasi", "Fondasi akademis kuat berpredikat Magna Cum Laude dipadukan dengan sertifikasi industri.")

    # Left Card: Education & Honors
    edu_card = add_card(slide9, Inches(0.8), Inches(1.95), Inches(5.7), Inches(5.0), BG_CARD_WHITE, BORDER_SUBTLE)
    
    # Header strip
    edu_strip = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.95), Inches(5.7), Inches(0.1))
    edu_strip.fill.solid()
    edu_strip.fill.fore_color.rgb = BLUE_ACCENT
    edu_strip.line.fill.background()

    tb_edu = slide9.shapes.add_textbox(Inches(1.05), Inches(2.2), Inches(5.2), Inches(4.6))
    tf_edu = tb_edu.text_frame
    tf_edu.word_wrap = True
    tf_edu.margin_left = tf_edu.margin_top = 0

    p0 = tf_edu.paragraphs[0]
    p0.text = "PENDIDIKAN FORMAL"
    p0.font.name = FONT_FAMILY
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = BLUE_ACCENT

    p1 = tf_edu.add_paragraph()
    p1.text = "Universitas Budi Luhur Jakarta"
    p1.font.name = FONT_FAMILY
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = NAVY_PRIMARY

    p2 = tf_edu.add_paragraph()
    p2.text = "Sarjana Sistem Informasi (S.Kom.) — Peminatan Data Science\nPeriode Studi: 2022 – 2026"
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_MUTED

    # Amber honor box
    h_box = tf_edu.add_paragraph()
    h_box.text = "\n⭐ PREDIKAT: MAGNA CUM LAUDE (IPK 3.88 / 4.00)"
    h_box.font.name = FONT_FAMILY
    h_box.font.size = Pt(12)
    h_box.font.bold = True
    h_box.font.color.rgb = RGBColor(180, 83, 9)

    edu_details = [
        "Fokus Kurikulum: Algoritma & Pemrograman, Database Relasional, Data Mining, Machine Learning, Statistika Komputasi, Business Intelligence.",
        "Publikasi Skripsi: Komparasi SVM, Naive Bayes dan KNN dalam Analisis Sentimen Respons Netizen YouTube Kasus Korupsi (Nilai Sidang: A).",
        "Aktivitas: Aktif dalam riset analitika data terapan dan pengembangan solusi otomasi kecerdasan buatan."
    ]
    for pt in edu_details:
        p_pt = tf_edu.add_paragraph()
        p_pt.text = f"• {pt}"
        p_pt.font.name = FONT_FAMILY
        p_pt.font.size = Pt(9.5)
        p_pt.font.color.rgb = TEXT_MAIN

    # Right Card: Professional Certifications & Continuous Upskilling
    cert_card = add_card(slide9, Inches(6.8), Inches(1.95), Inches(5.7), Inches(5.0), BG_CARD_WHITE, BORDER_SUBTLE)
    
    cert_strip = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.95), Inches(5.7), Inches(0.1))
    cert_strip.fill.solid()
    cert_strip.fill.fore_color.rgb = CYAN_ACCENT
    cert_strip.line.fill.background()

    tb_cert = slide9.shapes.add_textbox(Inches(7.05), Inches(2.2), Inches(5.2), Inches(4.6))
    tf_cert = tb_cert.text_frame
    tf_cert.word_wrap = True
    tf_cert.margin_left = tf_cert.margin_top = 0

    p_ct = tf_cert.paragraphs[0]
    p_ct.text = "SERTIFIKASI & KOMPETENSI INDUSTRI"
    p_ct.font.name = FONT_FAMILY
    p_ct.font.size = Pt(11)
    p_ct.font.bold = True
    p_ct.font.color.rgb = BLUE_ACCENT

    certs = [
        ("RevoU Data Analytics Specialization", "Business Analytics, SQL Querying, Tableau Storytelling, and Data-Driven Strategy formulation."),
        ("DQLab: Python for Data Science & Viz", "Data wrangling with Pandas/NumPy, exploratory data analysis, and statistical visualization."),
        ("Google Antigravity Agentic AI Mastery", "Autonomous multi-agent orchestration, LLM prompt engineering, and end-to-end automated pipelines."),
        ("Database Relasional & SQL Mastery", "Skema relational database design, CTEs, Window Functions, and query optimization.")
    ]

    for c_name, c_desc in certs:
        p_c = tf_cert.add_paragraph()
        p_c.text = f"✓ {c_name}"
        p_c.font.name = FONT_FAMILY
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = NAVY_PRIMARY

        p_cd = tf_cert.add_paragraph()
        p_cd.text = c_desc
        p_cd.font.name = FONT_FAMILY
        p_cd.font.size = Pt(9.5)
        p_cd.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 10: NILAI TAMBAH, TARGET PERAN & KONTAK PROFESIONAL
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_background(slide10, BG_LIGHT_SLATE)
    add_header(slide10, "LET'S COLLABORATE", "Nilai Tambah, Target Peran & Kontak Profesional", "Siap berkontribusi langsung memberikan dampak bisnis nyata melalui analitika data dan otomasi AI.")

    # Left Section: Value Proposition & Target Roles
    left_col_card = add_card(slide10, Inches(0.8), Inches(1.95), Inches(6.0), Inches(5.0), BG_CARD_WHITE, BORDER_SUBTLE)
    
    tb_vp = slide10.shapes.add_textbox(Inches(1.05), Inches(2.15), Inches(5.5), Inches(4.6))
    tf_vp = tb_vp.text_frame
    tf_vp.word_wrap = True
    tf_vp.margin_left = tf_vp.margin_top = 0

    p_vph = tf_vp.paragraphs[0]
    p_vph.text = "APA YANG SAYA HADIRKAN BAGI TIM ANDA?"
    p_vph.font.name = FONT_FAMILY
    p_vph.font.size = Pt(11)
    p_vph.font.bold = True
    p_vph.font.color.rgb = BLUE_ACCENT

    vp_items = [
        ("Kombinasi Teknis & Bisnis", "Tidak hanya mengolah angka, tetapi mampu mengartikulasikan insight analitika menjadi rekomendasi keputusan strategis bagi manajemen."),
        ("Riset Machine Learning Teruji", "Terbukti mampu merancang pipeline ML NLP akurasi tinggi (88.4%) dengan metodologi evaluasi yang valid dan bebas bias."),
        ("Adopsi Teknologi Agentic AI Terdepan", "Membawa kapabilitas otomatisasi masa depan untuk menghemat biaya operasional dan mempercepat time-to-delivery proyek.")
    ]

    for v_title, v_desc in vp_items:
        p_v = tf_vp.add_paragraph()
        p_v.text = f"❖ {v_title}: "
        p_v.font.name = FONT_FAMILY
        p_v.font.size = Pt(10)
        p_v.font.bold = True
        p_v.font.color.rgb = NAVY_PRIMARY
        r = p_v.add_run()
        r.text = v_desc
        r.font.bold = False
        r.font.color.rgb = TEXT_MUTED

    # Target Roles Box
    p_trh = tf_vp.add_paragraph()
    p_trh.text = "\nTARGET PERAN YANG SESUAI:"
    p_trh.font.name = FONT_FAMILY
    p_trh.font.size = Pt(10.5)
    p_trh.font.bold = True
    p_trh.font.color.rgb = NAVY_PRIMARY

    roles = ["Data Analyst", "Business Intelligence Analyst", "Junior Data Scientist", "Agentic AI Specialist"]
    for r_item in roles:
        p_r = tf_vp.add_paragraph()
        p_r.text = f"• {r_item}"
        p_r.font.name = FONT_FAMILY
        p_r.font.size = Pt(9.5)
        p_r.font.bold = True
        p_r.font.color.rgb = BLUE_ACCENT

    # Right Section: Direct Contact Cards
    right_contact_card = add_card(slide10, Inches(7.1), Inches(1.95), Inches(5.4), Inches(5.0), NAVY_PRIMARY, border_color=None)
    
    tb_rc = slide10.shapes.add_textbox(Inches(7.35), Inches(2.2), Inches(4.9), Inches(4.5))
    tf_rc = tb_rc.text_frame
    tf_rc.word_wrap = True
    tf_rc.margin_left = tf_rc.margin_top = 0

    p_ch = tf_rc.paragraphs[0]
    p_ch.text = "HUBUNGI SAYA LANGSUNG"
    p_ch.font.name = FONT_FAMILY
    p_ch.font.size = Pt(12)
    p_ch.font.bold = True
    p_ch.font.color.rgb = CYAN_ACCENT

    p_cs = tf_rc.add_paragraph()
    p_cs.text = "Terbuka untuk peluang kerja Full-time, Kontrak, atau Kolaborasi Proyek Strategis."
    p_cs.font.name = FONT_FAMILY
    p_cs.font.size = Pt(10)
    p_cs.font.color.rgb = RGBColor(226, 232, 240)

    contact_items = [
        ("Email Utama", "khairulraihan1204@gmail.com"),
        ("WhatsApp / Seluler", "+62 822-1081-3652"),
        ("Profil LinkedIn", "linkedin.com/in/khairulraihan"),
        ("Repositori GitHub", "github.com/KhairulRaihan"),
        ("Web Portofolio", "khairulraihan.github.io"),
        ("Lokasi Kerja", "Jakarta, Indonesia (On-site / Hybrid / Remote)")
    ]

    for lbl, val in contact_items:
        p_lbl = tf_rc.add_paragraph()
        p_lbl.text = f"{lbl.upper()}:"
        p_lbl.font.name = FONT_FAMILY
        p_lbl.font.size = Pt(8.5)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = CYAN_ACCENT

        p_val = tf_rc.add_paragraph()
        p_val.text = val
        p_val.font.name = FONT_FAMILY
        p_val.font.size = Pt(10.5)
        p_val.font.color.rgb = TEXT_WHITE

    # Save to both requested filenames
    out_files = [
        "Biru Putih Minimalis Modern Profesional Portofolio Presentasi.pptx",
        "Khairul_Raihan_Hidayat_Portfolio_Deck.pptx",
        "assets/docs/Khairul_Raihan_Hidayat_Portfolio_Deck.pptx"
    ]

    for fpath in out_files:
        os.makedirs(os.path.dirname(fpath) if os.path.dirname(fpath) else ".", exist_ok=True)
        prs.save(fpath)
        print(f"Successfully generated: {fpath}")

if __name__ == "__main__":
    build_presentation()
