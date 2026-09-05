"""
Script to generate a 10-slide executive presentation deck (.pptx)
for Khairul Raihan Hidayat's Portfolio.
"""

import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette - Obsidian & Slate Modern Aesthetics
    BG_COLOR = RGBColor(9, 11, 16)
    CARD_BG = RGBColor(18, 24, 34)
    CARD_BORDER = RGBColor(35, 47, 66)
    TEXT_MAIN = RGBColor(248, 250, 252)
    TEXT_MUTED = RGBColor(148, 163, 184)
    TEXT_DIM = RGBColor(100, 116, 139)
    ACCENT_BLUE = RGBColor(59, 130, 246)
    ACCENT_CYAN = RGBColor(6, 182, 212)
    ACCENT_EMERALD = RGBColor(16, 185, 129)
    ACCENT_PURPLE = RGBColor(139, 92, 246)
    ACCENT_AMBER = RGBColor(245, 158, 11)

    FONT_TITLE = "Outfit"
    FONT_BODY = "Plus Jakarta Sans"
    FONT_MONO = "JetBrains Mono"

    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_COLOR
        bg_shape.line.fill.background()
        return bg_shape

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.2)
        else:
            card.line.fill.background()
        card.adjustments[0] = 0.04
        return card

    def add_header(slide, badge_text, title_text, subtitle_text=None):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        # Badge
        p_badge = tf.paragraphs[0]
        p_badge.text = badge_text.upper()
        p_badge.font.name = FONT_MONO
        p_badge.font.size = Pt(9.5)
        p_badge.font.bold = True
        p_badge.font.color.rgb = ACCENT_CYAN
        p_badge.space_after = Pt(4)

        # Title
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_MAIN

        # Subtitle
        if subtitle_text:
            p_sub = tf.add_paragraph()
            p_sub.text = subtitle_text
            p_sub.font.name = FONT_BODY
            p_sub.font.size = Pt(11)
            p_sub.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 1: COVER SLIDE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    # Accent bar on left
    accent_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.08), Inches(3.8))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_BLUE
    accent_bar.line.fill.background()

    # Cover text frame
    s1_tb = s1.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(7.5), Inches(4.0))
    s1_tf = s1_tb.text_frame
    s1_tf.word_wrap = True
    s1_tf.margin_left = s1_tf.margin_top = s1_tf.margin_right = s1_tf.margin_bottom = 0

    p_badge = s1_tf.paragraphs[0]
    p_badge.text = "EXECUTIVE PORTFOLIO • DATA & AI ENGINEERING"
    p_badge.font.name = FONT_MONO
    p_badge.font.size = Pt(10)
    p_badge.font.bold = True
    p_badge.font.color.rgb = ACCENT_CYAN
    p_badge.space_after = Pt(10)

    p_name = s1_tf.add_paragraph()
    p_name.text = "Khairul Raihan Hidayat"
    p_name.font.name = FONT_TITLE
    p_name.font.size = Pt(36)
    p_name.font.bold = True
    p_name.font.color.rgb = TEXT_MAIN
    p_name.space_after = Pt(8)

    p_role = s1_tf.add_paragraph()
    p_role.text = "Data Science, Business Intelligence & Agentic AI Specialist"
    p_role.font.name = FONT_BODY
    p_role.font.size = Pt(16)
    p_role.font.bold = True
    p_role.font.color.rgb = ACCENT_BLUE
    p_role.space_after = Pt(16)

    p_bio = s1_tf.add_paragraph()
    p_bio.text = (
        "Lulusan Sarjana Sistem Informasi (IPK 3.88 / 4.00, Magna Cum Laude) dari Universitas Budi Luhur. "
        "Fokus pada siklus analitik komprehensif: pemodelan Machine Learning (NLP), perancangan dasbor eksekutif "
        "di Tableau & Looker Studio, serta orchestrasi alur kerja otonom berbasis Agentic AI (Antigravity)."
    )
    p_bio.font.name = FONT_BODY
    p_bio.font.size = Pt(11.5)
    p_bio.font.color.rgb = TEXT_MUTED
    p_bio.space_after = Pt(20)

    p_meta = s1_tf.add_paragraph()
    p_meta.text = "Tangerang, Banten • khairulraihan617@gmail.com • +62 898-9518-334"
    p_meta.font.name = FONT_MONO
    p_meta.font.size = Pt(9.5)
    p_meta.font.color.rgb = TEXT_DIM

    # Profile Image on Right Card
    add_card(s1, Inches(9.0), Inches(1.5), Inches(3.5), Inches(4.5))
    if os.path.exists("assets/images/profile.jpg"):
        s1.shapes.add_picture("assets/images/profile.jpg", Inches(9.2), Inches(1.7), Inches(3.1), Inches(4.1))

    # =========================================================================
    # SLIDE 2: EXECUTIVE SUMMARY & PROFESSIONAL PROFILE
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "Profil Profesional & Landasan Kerja", "Executive Summary", "Sintesis kualifikasi akademis, etos analitis, dan proposisi nilai bisnis.")

    # 3 Column Cards
    col_w = Inches(3.68)
    col_h = Inches(4.8)
    gap = Inches(0.32)
    start_x = Inches(0.8)
    start_y = Inches(1.85)

    # Card 1: Academic Rigor
    add_card(s2, start_x, start_y, col_w, col_h)
    tb1 = s2.shapes.add_textbox(start_x + Inches(0.25), start_y + Inches(0.25), col_w - Inches(0.5), col_h - Inches(0.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "01 • AKADEMIS & STATISTIK"
    p.font.name = FONT_MONO
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(6)

    p = tf1.add_paragraph()
    p.text = "Universitas Budi Luhur"
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(4)

    p = tf1.add_paragraph()
    p.text = "S.Kom Sistem Informasi (Data Science)\nIPK 3.88 / 4.00 (Magna Cum Laude)"
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(12)

    p = tf1.add_paragraph()
    p.text = (
        "• Penguasaan fondasi logika algoritma, kompleksitas komputasi, dan struktur data.\n"
        "• Pemodelan matematika & statistika prediktif: Regresi, Klasifikasi Teks, dan K-Means.\n"
        "• Lulusan terbaik dengan portofolio riset komparasi algoritma NLP terpublikasi."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED

    # Card 2: End-to-End Analytics
    add_card(s2, start_x + col_w + gap, start_y, col_w, col_h)
    tb2 = s2.shapes.add_textbox(start_x + col_w + gap + Inches(0.25), start_y + Inches(0.25), col_w - Inches(0.5), col_h - Inches(0.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "02 • BUSINESS INTELLIGENCE"
    p.font.name = FONT_MONO
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p.space_after = Pt(6)

    p = tf2.add_paragraph()
    p.text = "Decision-Making BI"
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(4)

    p = tf2.add_paragraph()
    p.text = "Tableau, SQL, Looker & Excel Lanjutan\nMulti-Page Executive Dashboards"
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p.space_after = Pt(12)

    p = tf2.add_paragraph()
    p.text = (
        "• Keahlian mendalam dalam Calculated Fields, agregasi dinamis, dan parameterisasi visual.\n"
        "• Transformasi data mentah puluhan ribu baris menjadi KPI yang actionable.\n"
        "• Visualisasi geospasial, matriks korelasi kinerja, dan pelaporan distribusi kompensasi."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED

    # Card 3: Modern Agentic AI
    add_card(s2, start_x + (col_w + gap) * 2, start_y, col_w, col_h)
    tb3 = s2.shapes.add_textbox(start_x + (col_w + gap) * 2 + Inches(0.25), start_y + Inches(0.25), col_w - Inches(0.5), col_h - Inches(0.5))
    tf3 = tb3.text_frame
    tf3.word_wrap = True

    p = tf3.paragraphs[0]
    p.text = "03 • AGENTIC AI AUTOMATION"
    p.font.name = FONT_MONO
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    p.space_after = Pt(6)

    p = tf3.add_paragraph()
    p.text = "Alur Kerja Otonom AI"
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(4)

    p = tf3.add_paragraph()
    p.text = "Antigravity, Scripting & Multi-Tools\nEfisiensi Kurasi Waktu ~70%"
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    p.space_after = Pt(12)

    p = tf3.add_paragraph()
    p.text = (
        "• Desain autonomous workflow untuk membedah video panjang & transkrip secara cepat.\n"
        "• Deteksi hook berpotensi retensi tinggi secara otomatis untuk kampanye konten.\n"
        "• Mengintegrasikan AI reasoning ke dalam eksekusi tugas teknis harian."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 3: SKILLS TAXONOMY & CORE COMPETENCIES
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "Taksonomi Keahlian Sesuai CV", "Core Competencies & Tools Matrix", "Klasifikasi keahlian teknis (Hard Skills) dan kompetensi interpersonal (Soft Skills).")

    # 4 Quadrants for Skills
    quad_w = Inches(5.7)
    quad_h = Inches(2.25)

    # Q1: Programming & Database
    add_card(s3, Inches(0.8), Inches(1.85), quad_w, quad_h)
    q1_tb = s3.shapes.add_textbox(Inches(1.0), Inches(1.95), quad_w - Inches(0.4), quad_h - Inches(0.2))
    q1_tf = q1_tb.text_frame
    q1_tf.word_wrap = True
    p = q1_tf.paragraphs[0]
    p.text = "PEMROGRAMAN & BASIS DATA"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(6)

    p = q1_tf.add_paragraph()
    p.text = "• Python: Pandas, NumPy, Scikit-Learn, NLTK, Streamlit, Matplotlib\n" \
             "• SQL: Relational Schema, DDL/DML, Complex Joins, Aggregation\n" \
             "• Tools: VS Code, Git / GitHub Version Control, Jupyter Notebook"
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED

    # Q2: BI & Visualization
    add_card(s3, Inches(6.8), Inches(1.85), quad_w, quad_h)
    q2_tb = s3.shapes.add_textbox(Inches(7.0), Inches(1.95), quad_w - Inches(0.4), quad_h - Inches(0.2))
    q2_tf = q2_tb.text_frame
    q2_tf.word_wrap = True
    p = q2_tf.paragraphs[0]
    p.text = "VISUALISASI DATA & BUSINESS INTELLIGENCE"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p.space_after = Pt(6)

    p = q2_tf.add_paragraph()
    p.text = "• Tableau Desktop & Tableau Public: Calculated Fields, Parameter, Dashboard multi-page\n" \
             "• Google Looker Studio: Pelaporan KPI real-time & dasbor bisnis eksekutif\n" \
             "• Microsoft Excel Lanjutan: Pivot Tables, Charts, XLOOKUP, Data Cleansing"
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED

    # Q3: Data Science & ML
    add_card(s3, Inches(0.8), Inches(4.35), quad_w, quad_h)
    q3_tb = s3.shapes.add_textbox(Inches(1.0), Inches(4.45), quad_w - Inches(0.4), quad_h - Inches(0.2))
    q3_tf = q3_tb.text_frame
    q3_tf.word_wrap = True
    p = q3_tf.paragraphs[0]
    p.text = "DATA SCIENCE, MACHINE LEARNING & NLP"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(6)

    p = q3_tf.add_paragraph()
    p.text = "• Natural Language Processing (NLP): Tokenizing, Normalisasi Kata Slang, InSet Lexicon\n" \
             "• Supervised Learning: Support Vector Machine (SVM/SVC), Multinomial Naive Bayes\n" \
             "• Unsupervised Learning: K-Means Clustering, Customer Profiling, RapidMiner"
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED

    # Q4: Agentic AI & Creative
    add_card(s3, Inches(6.8), Inches(4.35), quad_w, quad_h)
    q4_tb = s3.shapes.add_textbox(Inches(7.0), Inches(4.45), quad_w - Inches(0.4), quad_h - Inches(0.2))
    q4_tf = q4_tb.text_frame
    q4_tf.word_wrap = True
    p = q4_tf.paragraphs[0]
    p.text = "AGENTIC AI, TOOLS & SOFT SKILLS"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    p.space_after = Pt(6)

    p = q4_tf.add_paragraph()
    p.text = "• Agentic AI: Antigravity Autonomous Workflow, Task Orchestration, Prompt Craft\n" \
             "• Desain & Media: Figma, Canva, CapCut (9:16 Shorts Framing, Dynamic Captions)\n" \
             "• Karakter Kerja: Critical Thinking, Analytical Synthesis, Communication & Teamwork"
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 4: PROJECT 1 - TUGAS AKHIR / SKRIPSI NLP
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "Proyek Unggulan 01 • Tugas Akhir / Skripsi", "Analisis Sentimen Respons Netizen (YouTube)", "Pemodelan komparatif SVM vs Naive Bayes pada data komentar YouTube berbahasa Indonesia.")

    # Left: Details Card
    add_card(s4, Inches(0.8), Inches(1.85), Inches(6.8), Inches(4.9))
    s4_tb = s4.shapes.add_textbox(Inches(1.05), Inches(2.05), Inches(6.3), Inches(4.5))
    s4_tf = s4_tb.text_frame
    s4_tf.word_wrap = True

    p = s4_tf.paragraphs[0]
    p.text = "RINGKASAN & LATAR BELAKANG"
    p.font.name = FONT_MONO
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(6)

    p = s4_tf.add_paragraph()
    p.text = (
        "Bahasa ulasan netizen di media sosial sarat dengan bahasa tidak baku, singkatan, dan slang daerah. "
        "Riset ini merancang pipeline Machine Learning terintegrasi untuk mengekstraksi dan mengklasifikasikan "
        "persepsi publik secara otomatis berskala besar."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = s4_tf.add_paragraph()
    p.text = "FITUR KUNCI & METODOLOGI RISET"
    p.font.name = FONT_MONO
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(6)

    p = s4_tf.add_paragraph()
    p.text = (
        "1. Ekstraksi Data: Mengambil belasan ribu komentar YouTube via YouTube Data API v3.\n"
        "2. Normalisasi Teks: Kamus lokal kata baku (kamuskatabaku.xlsx) + leksikon InSet.\n"
        "3. Pemodelan: Komparasi algoritma Support Vector Machine (SVC) dan Naive Bayes.\n"
        "4. Deployment: Aplikasi web interaktif berbasis Streamlit untuk testing real-time."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = s4_tf.add_paragraph()
    p.text = "HASIL CAPAIAN AKURASI"
    p.font.name = FONT_MONO
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p.space_after = Pt(4)

    p = s4_tf.add_paragraph()
    p.text = "• Model SVM: Akurasi 92.4% (Precision & Recall optimal)\n• Model Naive Bayes: Akurasi 91.6%\n• Tech Stack: Python, Streamlit, Scikit-Learn, NLTK, Pandas"
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MAIN

    # Right: Metric Badges & Screenshot
    add_card(s4, Inches(7.8), Inches(1.85), Inches(4.7), Inches(4.9))
    if os.path.exists("assets/images/project-nlp.jpg"):
        s4.shapes.add_picture("assets/images/project-nlp.jpg", Inches(8.0), Inches(2.05), Inches(4.3), Inches(2.6))

    # Metric boxes below image
    m_box = s4.shapes.add_textbox(Inches(8.0), Inches(4.8), Inches(4.3), Inches(1.8))
    m_tf = m_box.text_frame
    m_tf.word_wrap = True
    p = m_tf.paragraphs[0]
    p.text = "METRIK UTAMA PROYEK"
    p.font.name = FONT_MONO
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p = m_tf.add_paragraph()
    p.text = "92.4% Akurasi SVM  |  12,500+ Komentar  |  91.6% Naive Bayes"
    p.font.name = FONT_TITLE
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p.space_after = Pt(4)

    p = m_tf.add_paragraph()
    p.text = "Teruji pada dataset riil komentar edukasi publik YouTube dengan pipeline preprocessing lengkap."
    p.font.name = FONT_BODY
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 5: NLP PIPELINE & STREAMLIT ARCHITECTURE
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Arsitektur Pipeline NLP • Skripsi", "Alur Data & Sistem Klasifikasi Sentimen", "Tahapan pengolahan data mentah dari YouTube API hingga pengujian interaktif di Streamlit.")

    # 4 Pipeline Stage Cards
    step_w = Inches(2.75)
    step_h = Inches(4.8)
    gap_s = Inches(0.24)
    sx = Inches(0.8)
    sy = Inches(1.85)

    stages = [
        ("TAHAP 01", "Ekstraksi & Scraping", ACCENT_CYAN,
         "• YouTube Data API v3\n• Pengambilan komentar terstruktur (Author, Comment, Likes, Date)\n• Penyimpanan dataset_mentah.csv\n• 12.500+ sampel data ulasan netizen."),
        ("TAHAP 02", "Text Preprocessing", ACCENT_BLUE,
         "• Case Folding & Cleansing teks\n• Kamus Kata Baku (kamuskatabaku.xlsx) untuk normalisasi slang\n• Stopword Removal bahasa Indonesia\n• Pembobotan Leksikon InSet."),
        ("TAHAP 03", "Feature & Training", ACCENT_PURPLE,
         "• CountVectorizer / TF-IDF\n• Split Train-Test (80:20)\n• Pelatihan SVM (Linear/RBF kernel)\n• Evaluasi Confusion Matrix & Classification Report (F1-score)."),
        ("TAHAP 04", "Streamlit Deployment", ACCENT_EMERALD,
         "• Dasbor visualisasi responsif (App.py)\n• Console uji kalimat langsung\n• Visualisasi WordCloud kata positif vs negatif\n• Grafik perbandingan performa algoritma.")
    ]

    for idx, (tag, title, color, desc) in enumerate(stages):
        x = sx + idx * (step_w + gap_s)
        add_card(s5, x, sy, step_w, step_h)
        tb = s5.shapes.add_textbox(x + Inches(0.18), sy + Inches(0.2), step_w - Inches(0.36), step_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = tag
        p.font.name = FONT_MONO
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(10)

        p = tf.add_paragraph()
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(9.8)
        p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 6: PROJECT 2 - TABLEAU HR ANALYTICS DASHBOARD
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "Proyek Unggulan 02 • Business Intelligence", "Human Resources (HR) Analytics Dashboard", "Perancangan dasbor eksekutif SDM dua halaman di Tableau menganalisis 8.950 catatan karyawan.")

    # Left: Details
    add_card(s6, Inches(0.8), Inches(1.85), Inches(5.8), Inches(4.9))
    s6_tb = s6.shapes.add_textbox(Inches(1.05), Inches(2.05), Inches(5.3), Inches(4.5))
    s6_tf = s6_tb.text_frame
    s6_tf.word_wrap = True

    p = s6_tf.paragraphs[0]
    p.text = "OVERVIEW PROYEK & KEBUTUHAN BISNIS"
    p.font.name = FONT_MONO
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p.space_after = Pt(6)

    p = s6_tf.add_paragraph()
    p.text = (
        "Membangun dasbor analitik SDM interaktif dua halaman (HR Overview & Employee Directory) "
        "berdasarkan simulasi dataset riil 8.950 karyawan untuk memonitor retensi, mengidentifikasi tren perputaran tenaga kerja, "
        "serta mengevaluasi kesetaraan kompensasi per departemen."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = s6_tf.add_paragraph()
    p.text = "METRIK KUNCI & INSIGHT BISNIS"
    p.font.name = FONT_MONO
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(6)

    p = s6_tf.add_paragraph()
    p.text = (
        "• Total Karyawan: 8.950 orang (7.984 Aktif / 89.2% retensi, 966 Terminasi)\n"
        "• Departemen Kunci: Operations (2.429 aktif), Sales (1.634 aktif), Customer Service, IT\n"
        "• Analisis Demografi: Rasio Gender (46% Pria, 54% Wanita), Sebaran Usia & Pendidikan\n"
        "• Geospasial: State Map sebaran tenaga kerja (HQ 70% vs Branch 30%)\n"
        "• Kompensasi & Kinerja: Scatter plot korelasi gaji vs usia & matriks rating performa."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED

    # Right: Dashboard Screenshot
    add_card(s6, Inches(6.8), Inches(1.85), Inches(5.7), Inches(4.9))
    if os.path.exists("assets/images/project-hr-analytics.png"):
        s6.shapes.add_picture("assets/images/project-hr-analytics.png", Inches(7.0), Inches(2.05), Inches(5.3), Inches(3.2))

    # Tableau Details note below image
    t_box = s6.shapes.add_textbox(Inches(7.0), Inches(5.4), Inches(5.3), Inches(1.2))
    t_tf = t_box.text_frame
    t_tf.word_wrap = True
    p = t_tf.paragraphs[0]
    p.text = "IMPLEMENTASI TEKNIS TABLEAU"
    p.font.name = FONT_MONO
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD

    p = t_tf.add_paragraph()
    p.text = "Mengembangkan Calculated Fields khusus, Interactive Action Filters, Dark Theme UI, dan Ekspor Laporan."
    p.font.name = FONT_BODY
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 7: PROJECT 3 - AGENTIC AI WORKFLOW SPECIALIST
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "Proyek Unggulan 03 • Agentic AI Automation", "AI Content Workflow & Analytics Specialist", "Otomasi kurasi konten video pendek menggunakan Antigravity dan alur kerja AI terorkestrasi.")

    # Left: Big Highlight Card
    add_card(s7, Inches(0.8), Inches(1.85), Inches(6.0), Inches(4.9))
    s7_tb = s7.shapes.add_textbox(Inches(1.05), Inches(2.05), Inches(5.5), Inches(4.5))
    s7_tf = s7_tb.text_frame
    s7_tf.word_wrap = True

    p = s7_tf.paragraphs[0]
    p.text = "TANTANGAN & MASALAH"
    p.font.name = FONT_MONO
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    p.space_after = Pt(6)

    p = s7_tf.add_paragraph()
    p.text = (
        "Menonton dan mengkurasi momen bernilai tinggi (high-retention hooks) dari video panjang berdurasi "
        "puluhan menit membutuhkan waktu kerja manual yang sangat besar, membatasi skalabilitas kampanye konten."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = s7_tf.add_paragraph()
    p.text = "SOLUSI BERBASIS AGENTIC AI"
    p.font.name = FONT_MONO
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(6)

    p = s7_tf.add_paragraph()
    p.text = (
        "• Mengembangkan workflow otonom dengan Antigravity untuk menganalisis transkrip dan konteks video.\n"
        "• Deteksi hook berpotensi engagement tinggi secara instan menggunakan structured prompts.\n"
        "• Ekstraksi klip 9:16 di CapCut dengan dynamic captions dan ritme visual teroptimasi.\n"
        "• Penyelarasan metadata YouTube (titles, tags, description) untuk memaksimalkan algoritma."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED

    # Right: 3 Metric Impact Cards
    r_x = Inches(7.1)
    r_w = Inches(5.4)
    add_card(s7, r_x, Inches(1.85), r_w, Inches(1.45))
    rc1 = s7.shapes.add_textbox(r_x + Inches(0.3), Inches(2.0), r_w - Inches(0.6), Inches(1.15))
    p = rc1.text_frame.paragraphs[0]
    p.text = "~70% EFISIENSI WAKTU KERJA"
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p = rc1.text_frame.add_paragraph()
    p.text = "Memangkas waktu kurasi manual dari 3 jam menjadi kurang dari 45 menit per video panjang."
    p.font.name = FONT_BODY
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_MUTED

    add_card(s7, r_x, Inches(3.55), r_w, Inches(1.45))
    rc2 = s7.shapes.add_textbox(r_x + Inches(0.3), Inches(3.7), r_w - Inches(0.6), Inches(1.15))
    p = rc2.text_frame.paragraphs[0]
    p.text = "9:16 SHORTS FORMATTING"
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p = rc2.text_frame.add_paragraph()
    p.text = "Pengemasan klip terstandarisasi untuk YouTube Shorts dengan dynamic subtitle & visual pacing."
    p.font.name = FONT_BODY
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_MUTED

    add_card(s7, r_x, Inches(5.25), r_w, Inches(1.5))
    rc3 = s7.shapes.add_textbox(r_x + Inches(0.3), Inches(5.4), r_w - Inches(0.6), Inches(1.2))
    p = rc3.text_frame.paragraphs[0]
    p.text = "100% AUTOMATED HOOK PIPELINE"
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    p = rc3.text_frame.add_paragraph()
    p.text = "Alur kerja agentic AI yang konsisten untuk memenuhi target campaign Content Reward berskala tinggi."
    p.font.name = FONT_BODY
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 8: PROJECT 4 - AMAZON SALES & COMMERCIAL BI
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Proyek Tambahan 04 • Commercial Analytics", "Amazon Sales & Revenue Performance Dashboard", "Analisis data transaksi penjualan regional Amazon untuk strategi margin laba dan rantai pasok.")

    # Left: Details
    add_card(s8, Inches(0.8), Inches(1.85), Inches(5.8), Inches(4.9))
    s8_tb = s8.shapes.add_textbox(Inches(1.05), Inches(2.05), Inches(5.3), Inches(4.5))
    s8_tf = s8_tb.text_frame
    s8_tf.word_wrap = True

    p = s8_tf.paragraphs[0]
    p.text = "LATAR BELAKANG & KEBUTUHAN BISNIS"
    p.font.name = FONT_MONO
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AMBER
    p.space_after = Pt(6)

    p = s8_tf.add_paragraph()
    p.text = (
        "Dataset penjualan e-commerce berskala regional membutuhkan sentralisasi visual agar eksekutif "
        "dapat mengevaluasi disparitas profit margin antar kategori produk serta performa komparatif antara "
        "saluran penjualan Online dan Offline."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = s8_tf.add_paragraph()
    p.text = "FITUR ANALITIK UTAMA"
    p.font.name = FONT_MONO
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(6)

    p = s8_tf.add_paragraph()
    p.text = (
        "• Agregasi Metrik Komersial: Total Revenue, Total Cost, Total Profit, Units Sold.\n"
        "• Waktu Kirim (Ship Days): Identifikasi hambatan logistik antar prioritas pesanan (Critical, High, Medium, Low).\n"
        "• Komparasi Saluran: Online vs Offline channel volume di berbagai benua pasar.\n"
        "• Tech Stack: Tableau Desktop, SQL Querying, Data Modeling, Excel."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED

    # Right: Dashboard visual
    add_card(s8, Inches(6.8), Inches(1.85), Inches(5.7), Inches(4.9))
    if os.path.exists("assets/images/project-tableau.jpg"):
        s8.shapes.add_picture("assets/images/project-tableau.jpg", Inches(7.0), Inches(2.05), Inches(5.3), Inches(3.2))

    a_box = s8.shapes.add_textbox(Inches(7.0), Inches(5.4), Inches(5.3), Inches(1.2))
    a_tf = a_box.text_frame
    a_tf.word_wrap = True
    p = a_tf.paragraphs[0]
    p.text = "IMPAK BISNIS & STRATEGI KOMERSIAL"
    p.font.name = FONT_MONO
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AMBER

    p = a_tf.add_paragraph()
    p.text = "Membantu manajer penjualan memetakan produk dengan ROI tertinggi dan memangkas jeda logistik pengiriman."
    p.font.name = FONT_BODY
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 9: CERTIFICATIONS & ACADEMIC ACHIEVEMENTS
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "Kredensial & Bahasa", "Sertifikasi & Kualifikasi Resmi", "Pengakuan kompetensi teknis dari institusi akademis dan platform teknologi terkemuka.")

    # 4 Cards for Certifications
    cert_w = Inches(5.7)
    cert_h = Inches(2.25)

    certs = [
        ("ALGORITHM CERTIFICATION", "Universitas Budi Luhur • April 2024", ACCENT_BLUE,
         "Penguasaan logika algoritma, kompleksitas waktu komputasi, dan struktur data fundamental."),
        ("INTRODUCTION TO DATA ANALYTICS", "RevoU Tech Academy • 2024", ACCENT_EMERALD,
         "Metodologi analisis data bisnis, exploratory data analysis (EDA), dan visualisasi data praktis."),
        ("PYTHON FUNDAMENTAL FOR DATA SCIENCE", "DQLab Academy • Desember 2023", ACCENT_CYAN,
         "Sintaks Python dasar, manipulasi array & dataframe (Pandas/NumPy), dan pemrosesan data ilmiah."),
        ("R FUNDAMENTAL FOR DATA SCIENCE", "DQLab Academy • Desember 2023", ACCENT_PURPLE,
         "Komputasi statistik, struktur dataframe, dan visualisasi data dasar dalam bahasa R.")
    ]

    positions = [
        (Inches(0.8), Inches(1.85)),
        (Inches(6.8), Inches(1.85)),
        (Inches(0.8), Inches(4.35)),
        (Inches(6.8), Inches(4.35))
    ]

    for (title, issuer, color, desc), (px, py) in zip(certs, positions):
        add_card(s9, px, py, cert_w, cert_h)
        tb = s9.shapes.add_textbox(px + Inches(0.25), py + Inches(0.2), cert_w - Inches(0.5), cert_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = issuer
        p.font.name = FONT_MONO
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 10: CONCLUSION & CONTACT INFORMATION
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)

    # Big Center Card
    add_card(s10, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))

    c_tb = s10.shapes.add_textbox(Inches(1.8), Inches(1.5), Inches(9.7), Inches(4.5))
    c_tf = c_tb.text_frame
    c_tf.word_wrap = True

    p = c_tf.paragraphs[0]
    p.text = "KOLABORASI & KESEMPATAN KARIR"
    p.font.name = FONT_MONO
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(8)

    p = c_tf.add_paragraph()
    p.text = "Siap Berkontribusi Nyata dalam Analitik Data & AI"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(14)

    p = c_tf.add_paragraph()
    p.text = (
        "Terbuka untuk posisi Data Scientist, Data Analyst, Business Intelligence Specialist, "
        "atau AI Workflow Automation Specialist. Membawa kombinasi ketajaman akademis (IPK 3.88), "
        "keterampilan visualisasi Tableau, pemodelan Machine Learning Python, dan adaptabilitas Agentic AI."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(24)

    # Contact Box
    p = c_tf.add_paragraph()
    p.text = "INFORMASI KONTAK RESMI:"
    p.font.name = FONT_MONO
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(6)

    p = c_tf.add_paragraph()
    p.text = "• Email: khairulraihan617@gmail.com\n" \
             "• WhatsApp / Telepon: +62 898-9518-334\n" \
             "• LinkedIn: linkedin.com/in/khairul-raihan-hidayat-0a4b62334/\n" \
             "• GitHub: github.com/KhairulRaihan\n" \
             "• Portofolio Web: https://khairulraihan.github.io/\n" \
             "• Lokasi: Kota Tangerang, Banten, Indonesia"
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN

    # Save presentation
    output_pptx = "Khairul_Raihan_Hidayat_Portfolio_Deck.pptx"
    prs.save(output_pptx)
    print(f"Presentation saved successfully to: {output_pptx}")
    return output_pptx

if __name__ == "__main__":
    build_presentation()
