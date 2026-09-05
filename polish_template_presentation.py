import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Paths
BACKUP_PATH = r"C:\Users\Khairul Raihan\Desktop\career\Biru Putih Minimalis Modern Profesional Portofolio Presentasi_BACKUP.pptx"
OUTPUT_PATH = r"C:\Users\Khairul Raihan\Desktop\career\Biru Putih Minimalis Modern Profesional Portofolio Presentasi.pptx"
DOCS_PPTX_PATH = r"C:\Users\Khairul Raihan\Desktop\career\assets\docs\Biru_Putih_Portofolio_Presentasi.pptx"
CHART_PATH = r"C:\Users\Khairul Raihan\Desktop\career\slide_7_nlp_chart.png"

prs = pptx.Presentation(BACKUP_PATH)

NAVY = RGBColor(15, 36, 55)         # #0F2437
WHITE = RGBColor(255, 255, 255)     # #FFFFFF
STEEL = RGBColor(46, 92, 124)       # #2E5C7C
LIGHT_BLUE = RGBColor(200, 225, 245)# #C8E1F5

def set_text_exact(shape, text, font_name="DM Sans", font_size=Pt(20), bold=False, color=NAVY, alignment=None, zero_margin=False):
    """Safely updates a text box and cleanly strips any stale paragraphs from the template."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    
    if zero_margin:
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
    lines = text.split("\n")
    
    # Cleanly remove all existing paragraphs except the first one
    p0 = tf.paragraphs[0]
    p0.text = ""
    p_elems = list(tf._txBody.p_lst)
    for p_elem in p_elems[1:]:
        tf._txBody.remove(p_elem)
        
    # Write new lines
    for i, line in enumerate(lines):
        p = p0 if i == 0 else tf.add_paragraph()
        if alignment is not None:
            p.alignment = alignment
        p.text = ""
        if line:
            run = p.add_run()
            run.text = line
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = color

print("--> Slide 1: Cover")
s1 = prs.slides[0]
for s in s1.shapes:
    if s.shape_id == 21: # '2035'
        set_text_exact(s, "2024", font_name="DM Sans Bold", font_size=Pt(36), bold=True, color=NAVY)
    elif s.shape_id == 23: # Name
        s.top = Inches(5.15)
        s.width = Inches(9.8)
        set_text_exact(s, "Khairul Raihan Hidayat, S.Kom.", font_name="DM Sans Bold", font_size=Pt(42), bold=True, color=NAVY)
    elif s.shape_id == 22: # Role
        s.top = Inches(6.65)
        s.width = Inches(9.0)
        set_text_exact(s, "Data Analyst & Agentic AI Specialist", font_name="DM Sans", font_size=Pt(26), bold=False, color=STEEL)

print("--> Slide 2: Daftar Isi")
s2 = prs.slides[1]
for s in s2.shapes:
    if s.shape_id == 14:
        set_text_exact(s, "Ringkasan Profil", font_name="DM Sans", font_size=Pt(38), bold=False, color=WHITE)
    elif s.shape_id == 15:
        set_text_exact(s, "Latar Belakang Pendidikan", font_name="DM Sans", font_size=Pt(38), bold=False, color=WHITE)
    elif s.shape_id == 16:
        set_text_exact(s, "Keahlian Teknis & Alat Kerja", font_name="DM Sans", font_size=Pt(38), bold=False, color=WHITE)
    elif s.shape_id == 17:
        set_text_exact(s, "Studi Kasus & Riset Skripsi", font_name="DM Sans", font_size=Pt(38), bold=False, color=WHITE)
    elif s.shape_id == 18:
        set_text_exact(s, "Dampak, Metodologi & Kontak", font_name="DM Sans", font_size=Pt(38), bold=False, color=WHITE)

print("--> Slide 3: Ringkasan Profil")
s3 = prs.slides[2]
for s in s3.shapes:
    if s.shape_id == 23:
        summary_text = (
            "Lulusan Sarjana Komputer (S.Kom) Sistem Informasi konsentrasi Data Science dari Universitas Budi Luhur "
            "dengan predikat Magna Cum Laude (IPK 3.88/4.00).\n\n"
            "Memiliki spesialisasi dalam analisis data eksploratif, pemodelan prediktif Machine Learning (NLP), "
            "visualisasi Business Intelligence dengan Tableau, serta automasi alur kerja cerdas berbasis Agentic AI.\n\n"
            "Berpengalaman mengolah dataset skala besar, membangun model klasifikasi akurasi tinggi (92%+), "
            "dan merancang dasbor analitik interaktif untuk mendukung pengambilan keputusan bisnis yang presisi dan terukur."
        )
        set_text_exact(s, summary_text, font_name="DM Sans", font_size=Pt(20), bold=False, color=NAVY)

print("--> Slide 4: Latar Belakang")
s4 = prs.slides[3]
for s in s4.shapes:
    if s.shape_id == 28: # Card 1 Year
        set_text_exact(s, "2020 - 2024", font_name="DM Sans Bold", font_size=Pt(25), bold=True, color=NAVY)
    elif s.shape_id == 33: # Degree title
        set_text_exact(s, "Sarjana\nSistem Informasi", font_name="DM Sans Bold", font_size=Pt(28), bold=True, color=NAVY)
    elif s.shape_id == 34: # University & GPA
        set_text_exact(s, "Universitas Budi Luhur\nIPK 3.88 / 4.00 (Magna Cum Laude)", font_name="DM Sans", font_size=Pt(22), bold=False, color=STEEL)
    elif s.shape_id == 25: # Card 2 Year
        set_text_exact(s, "2024", font_name="DM Sans Bold", font_size=Pt(25), bold=True, color=NAVY)
    elif s.shape_id == 38: # Card 2 Course
        set_text_exact(s, "Mini Course\nData Analytics", font_name="DM Sans Bold", font_size=Pt(28), bold=True, color=NAVY)
    elif s.shape_id == 39: # Card 2 Academy
        set_text_exact(s, "RevoU", font_name="DM Sans", font_size=Pt(24), bold=False, color=STEEL)
    elif s.shape_id == 37: # Card 3 Year
        set_text_exact(s, "2024", font_name="DM Sans Bold", font_size=Pt(25), bold=True, color=NAVY)
    elif s.shape_id == 40: # Card 3 Certification
        set_text_exact(s, "Python Data Science &\nVisualization Mastery", font_name="DM Sans Bold", font_size=Pt(26), bold=True, color=NAVY)
    elif s.shape_id == 41: # Card 3 Academy
        set_text_exact(s, "DQLab Academy", font_name="DM Sans", font_size=Pt(24), bold=False, color=STEEL)
    elif s.shape_id == 29: # Stray shape
        set_text_exact(s, "", font_size=Pt(1))

print("--> Slide 5: Keahlian Teknis")
s5 = prs.slides[4]
for s in s5.shapes:
    # Card 1 (Top Left)
    if s.shape_id == 28:
        set_text_exact(s, "Bahasa Pemrograman & Data Science", font_name="DM Sans Bold", font_size=Pt(21), bold=True, color=WHITE)
    elif s.shape_id == 26:
        set_text_exact(s, "Python (Pandas, NumPy, Scikit-Learn, Streamlit), data wrangling, eksplorasi statistik (EDA), dan pemodelan prediktif machine learning terstruktur.", font_name="DM Sans", font_size=Pt(17), bold=False, color=WHITE)
    
    # Card 2 (Bottom Left)
    elif s.shape_id == 11:
        set_text_exact(s, "Kueri SQL & Manajemen Data", font_name="DM Sans Bold", font_size=Pt(21), bold=True, color=WHITE)
    elif s.shape_id == 12:
        set_text_exact(s, "Kueri SQL relasional (MySQL, PostgreSQL), agregasi data skala besar, operasi JOIN multi-tabel, optimasi query, dan ekstraksi dataset analitis.", font_name="DM Sans", font_size=Pt(17), bold=False, color=WHITE)
    
    # Card 3 (Top Right)
    elif s.shape_id == 27:
        set_text_exact(s, "Visualisasi Data & BI (Tableau)", font_name="DM Sans Bold", font_size=Pt(21), bold=True, color=WHITE)
    elif s.shape_id == 16:
        set_text_exact(s, "Tableau Desktop & Public, perancangan dasbor KPI interaktif, calculated fields, parameter dinamis, analisis tren bisnis, dan storytelling data eksekutif.", font_name="DM Sans", font_size=Pt(17), bold=False, color=WHITE)
    
    # Card 4 (Bottom Right)
    elif s.shape_id == 6:
        set_text_exact(s, "Agentic AI & Automasi", font_name="DM Sans Bold", font_size=Pt(21), bold=True, color=WHITE)
    elif s.shape_id == 7:
        set_text_exact(s, "Google Antigravity Agentic AI, rekayasa prompt LLM terstruktur, automasi batch processing, integrasi API, dan peningkatan efisiensi kerja hingga ~70%.", font_name="DM Sans", font_size=Pt(17), bold=False, color=WHITE)

print("--> Slide 6: Pengalaman Kerja & Riset")
s6 = prs.slides[5]
for s in s6.shapes:
    if s.shape_id == 26:
        set_text_exact(s, "Pengalaman Kerja & Riset", font_name="DM Sans Bold", font_size=Pt(60), bold=True, color=WHITE)
    elif s.shape_id == 29:
        s.width = Inches(8.5)
        set_text_exact(s, "AI CONTENT & WORKFLOW SPECIALIST", font_name="DM Sans Bold", font_size=Pt(24), bold=True, color=WHITE)
    elif s.shape_id == 30:
        set_text_exact(s, "Freelance / Remote | 2024 - Sekarang", font_name="DM Sans", font_size=Pt(22), bold=False, color=LIGHT_BLUE)
    elif s.shape_id == 27:
        set_text_exact(s, "Merancang pipeline automasi multi-agent menggunakan Google Antigravity & LLM untuk ekstraksi dan kurasi konten video format pendek (9:16).", font_name="DM Sans", font_size=Pt(18), bold=False, color=WHITE)
    elif s.shape_id == 28:
        set_text_exact(s, "Mengintegrasikan analisis metrik performa audiens dan automasi batch CapCut yang berhasil memangkas siklus produksi konten hingga ~70%.", font_name="DM Sans", font_size=Pt(18), bold=False, color=WHITE)
    
    elif s.shape_id == 31:
        s.width = Inches(8.5)
        set_text_exact(s, "PENELITI SKRIPSI DATA SCIENCE (NLP)", font_name="DM Sans Bold", font_size=Pt(24), bold=True, color=WHITE)
    elif s.shape_id == 32:
        set_text_exact(s, "Universitas Budi Luhur | 2023 - 2024", font_name="DM Sans", font_size=Pt(22), bold=False, color=LIGHT_BLUE)
    elif s.shape_id == 33:
        set_text_exact(s, "Melakukan riset komparatif algoritma SVM dan Multinomial Naive Bayes pada 12.500+ komentar YouTube otomotif mobil listrik (Wuling BinguoEV).", font_name="DM Sans", font_size=Pt(18), bold=False, color=WHITE)
    elif s.shape_id == 34:
        set_text_exact(s, "Mengembangkan pipeline NLP berbasis InSet Lexicon Bahasa Indonesia dan mendeploy aplikasi web analitik interaktif berbasis Streamlit.", font_name="DM Sans", font_size=Pt(18), bold=False, color=WHITE)
    elif s.shape_id == 5:
        set_text_exact(s, "", font_size=Pt(1))

print("--> Slide 7: Studi Kasus 1")
s7 = prs.slides[6]
for s in s7.shapes:
    if s.shape_id == 8:
        s.left = Inches(10.2)
    elif s.shape_id == 14:
        s.left = Inches(10.2)
        s.top = Inches(2.65)
        s.width = Inches(8.8)
        set_text_exact(s, "Analisis Sentimen (NLP Skripsi)", font_name="DM Sans", font_size=Pt(34), bold=False, color=STEEL)
    elif s.shape_id == 9:
        s.left = Inches(10.2)
        s.top = Inches(3.6)
        s.width = Inches(8.8)
        nlp_case_study = (
            "Riset Skripsi komparatif klasifikasi sentimen pada 12.500+ komentar YouTube terkait kendaraan listrik "
            "(Wuling BinguoEV) menggunakan algoritma Support Vector Machine (SVM) dan Multinomial Naive Bayes.\n\n"
            "Penelitian ini mengimplementasikan InSet Lexicon Bahasa Indonesia untuk penanganan teks informal, "
            "ekstraksi fitur TF-IDF, serta validasi silang (cross-validation).\n\n"
            "Hasil evaluasi membuktikan model SVM mencapai performa unggul dengan Akurasi 92.4%, Presisi 92.0%, dan "
            "Recall 92.8%, mengungguli Naive Bayes (91.6%). Seluruh pipeline dideploy ke dalam aplikasi web analitik interaktif berbasis Streamlit."
        )
        set_text_exact(s, nlp_case_study, font_name="DM Sans", font_size=Pt(19.5), bold=False, color=NAVY)

# Replace Picture 13 with custom chart
pic13 = None
for s in s7.shapes:
    if s.shape_id == 13 or s.name == "Picture 13":
        pic13 = s
        break

if pic13 and os.path.exists(CHART_PATH):
    p_left, p_top, p_width, p_height = pic13.left, pic13.top, pic13.width, pic13.height
    s7.shapes._spTree.remove(pic13._element)
    s7.shapes.add_picture(CHART_PATH, p_left, p_top, p_width, p_height)

print("--> Slide 8: Studi Kasus 2")
s8 = prs.slides[7]
for s in s8.shapes:
    if s.shape_id == 7:
        set_text_exact(s, "Dasbor Tableau HR & Komersial", font_name="DM Sans", font_size=Pt(42), bold=False, color=WHITE)
    elif s.shape_id == 5:
        hr_case_study = (
            "Perancangan dasbor Business Intelligence komprehensif menggunakan Tableau Desktop untuk menganalisis "
            "8.950 catatan karyawan lintas 7 departemen korporat.\n\n"
            "Dasbor ini menyajikan visualisasi KPI tingkat retensi (89.2%), analisis churn, scatter plot kompensasi vs performa, "
            "serta pemetaan sebaran demografis tenaga kerja untuk mendukung strategi retensi HR.\n\n"
            "Dilengkapi juga dengan analisis komersial Amazon Sales Dashboard yang memetakan profitabilitas regional, "
            "kinerja kategori produk, dan pemantauan SLA pengiriman guna mengoptimalkan margin penjualan e-commerce."
        )
        set_text_exact(s, hr_case_study, font_name="DM Sans", font_size=Pt(20.5), bold=False, color=WHITE)

print("--> Slide 9: Metodologi & Alur Kerja Data")
s9 = prs.slides[8]
for s in s9.shapes:
    if s.shape_id == 23:
        set_text_exact(s, "1. Pengumpulan & Ekstraksi Data", font_name="DM Sans Bold", font_size=Pt(21), bold=True, color=WHITE)
    elif s.shape_id == 24:
        set_text_exact(s, "Mengumpulkan 12.500+ komentar via YouTube Data API v3 dan mengintegrasikan 8.950 catatan data korporat terstruktur ke dalam pipeline analitis.", font_name="DM Sans", font_size=Pt(18.5), bold=False, color=WHITE)
    elif s.shape_id == 25:
        set_text_exact(s, "2. Pembersihan & Preprocessing Data", font_name="DM Sans Bold", font_size=Pt(21), bold=True, color=WHITE)
    elif s.shape_id == 26:
        set_text_exact(s, "Pembersihan teks (Case Folding, Tokenizing, Stopwords Removal), normalisasi kata slang, dan pembobotan polaritas sentimen berbasis InSet Lexicon.", font_name="DM Sans", font_size=Pt(18.5), bold=False, color=WHITE)
    elif s.shape_id == 27:
        set_text_exact(s, "3. Pemodelan, Evaluasi & Dasbor", font_name="DM Sans Bold", font_size=Pt(21), bold=True, color=WHITE)
    elif s.shape_id == 28:
        set_text_exact(s, "Ekstraksi fitur TF-IDF, pelatihan model SVM vs Naive Bayes (akurasi 92.4%), serta deployment dasbor interaktif berbasis Streamlit dan Tableau.", font_name="DM Sans", font_size=Pt(18.5), bold=False, color=WHITE)

print("--> Slide 10: Dampak & Angka Pencapaian")
s10 = prs.slides[9]
for s in s10.shapes:
    # Card Top Left (id 26, 27, 28)
    if s.shape_id == 26:
        s.width = Inches(2.0)
        set_text_exact(s, "92.4%", font_name="DM Sans Bold", font_size=Pt(38), bold=True, color=STEEL, zero_margin=True)
    elif s.shape_id == 27:
        s.width = Inches(4.5)
        set_text_exact(s, "Akurasi Model Machine Learning (SVM)", font_name="DM Sans Bold", font_size=Pt(21), bold=True, color=NAVY)
    elif s.shape_id == 28:
        set_text_exact(s, "Akurasi klasifikasi sentimen Skripsi mengungguli Naive Bayes (91.6%) dengan presisi 92.0% dan recall 92.8% pada 12.500+ data komentar YouTube.", font_name="DM Sans", font_size=Pt(18), bold=False, color=NAVY)
    
    # Card Top Right (id 19, 20, 21)
    elif s.shape_id == 19:
        s.width = Inches(2.1)
        set_text_exact(s, "8.950+", font_name="DM Sans Bold", font_size=Pt(38), bold=True, color=STEEL, zero_margin=True)
    elif s.shape_id == 20:
        s.width = Inches(4.5)
        set_text_exact(s, "Catatan Karyawan Dianalisis (Tableau)", font_name="DM Sans Bold", font_size=Pt(21), bold=True, color=NAVY)
    elif s.shape_id == 21:
        set_text_exact(s, "Pemodelan analitik retensi tenaga kerja lintas 7 departemen dan berbagai wilayah geografis pada dasbor Tableau Business Intelligence.", font_name="DM Sans", font_size=Pt(18), bold=False, color=NAVY)
    
    # Card Bottom Left (id 33, 34, 35)
    elif s.shape_id == 33:
        s.width = Inches(2.0)
        set_text_exact(s, "~70%", font_name="DM Sans Bold", font_size=Pt(38), bold=True, color=STEEL, zero_margin=True)
    elif s.shape_id == 34:
        s.width = Inches(2.5)
        set_text_exact(s, "Efisiensi Kerja\nKonten AI", font_name="DM Sans Bold", font_size=Pt(19), bold=True, color=NAVY)
    elif s.shape_id == 35:
        set_text_exact(s, "Pengurangan waktu kurasi dan produksi video format pendek (9:16) secara signifikan menggunakan pipeline automasi multi-agent AI.", font_name="DM Sans", font_size=Pt(18), bold=False, color=NAVY)
    
    # Card Bottom Right (id 12, 13, 14)
    elif s.shape_id == 12:
        s.width = Inches(2.0)
        set_text_exact(s, "3.88", font_name="DM Sans Bold", font_size=Pt(38), bold=True, color=STEEL, zero_margin=True)
    elif s.shape_id == 13:
        s.width = Inches(4.5)
        set_text_exact(s, "IPK Kelulusan (Magna Cum Laude)", font_name="DM Sans Bold", font_size=Pt(21), bold=True, color=NAVY)
    elif s.shape_id == 14:
        set_text_exact(s, "Predikat kehormatan Sarjana Komputer (S.Kom) Program Studi Sistem Informasi konsentrasi Data Science di Universitas Budi Luhur.", font_name="DM Sans", font_size=Pt(18), bold=False, color=NAVY)

print("--> Slide 11: Kontak & Penutup")
s11 = prs.slides[10]
for s in s11.shapes:
    if s.shape_id == 23:
        set_text_exact(s, "2024", font_name="DM Sans Bold", font_size=Pt(36), bold=True, color=WHITE)
    elif s.shape_id == 22:
        s.width = Inches(6.5)
        set_text_exact(s, "+62 898-9518-334", font_name="DM Sans", font_size=Pt(25), bold=False, color=WHITE)
    elif s.shape_id == 25:
        s.width = Inches(8.5)
        contact_text = "khairulraihan617@gmail.com\nlinkedin.com/in/khairul-raihan-hidayat\nkhairulraihan.github.io"
        set_text_exact(s, contact_text, font_name="DM Sans", font_size=Pt(21), bold=False, color=WHITE)



# Save presentations
prs.save(OUTPUT_PATH)
prs.save(DOCS_PPTX_PATH)
print("=== SAVED PPTX SUCCESSFULLY TO ROOT AND DOCS! ===")
